import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # 1. Initialize Presentation & Widescreen (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide layout

    # Colors (Light Theme based on Brand Logo & Guidelines)
    COLOR_BG = RGBColor(245, 247, 250)    # Soft Alabaster White
    COLOR_CARD = RGBColor(255, 255, 255)  # Pure White Card
    COLOR_GOLD = RGBColor(170, 132, 32)   # Brushed Brass Gold
    COLOR_BLUE = RGBColor(12, 125, 203)   # Logo Sky Blue
    COLOR_TEXT = RGBColor(13, 23, 39)     # Deep Navy Black Text
    COLOR_MUTED = RGBColor(92, 108, 132)  # Muted Slate Gray Text

    # Image Paths
    img_hero = "assets/bright_kitchen_hero.png"
    img_display = "assets/bright_appliances_display.png"
    logo_path = "assets/dreamhome_logo.jpg"

    # Helpers
    def apply_slide_background(slide, image_path=None):
        # 1. Set solid fallback background (light)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # 2. Add full screen background image if available
        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(image_path, 0, 0, Inches(13.333), Inches(7.5))

            # 3. Add full-screen light overlay for text contrast
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = COLOR_BG
            overlay.fill.transparency = 0.72  # 72% opacity white overlay
            overlay.line.fill.background()    # No border

    def add_slide_header(slide, title_text, category_text="INVESTOR & COLLABORATION PROFILE"):
        # Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(9.0), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(8.5)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_GOLD

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(9.0), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Georgia'
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT

        # Watermark Logo
        if os.path.exists(logo_path):
            logo_wm = slide.shapes.add_picture(logo_path, Inches(11.833), Inches(0.4), Inches(0.7), Inches(0.7))
            wm_frame = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.833), Inches(0.4), Inches(0.7), Inches(0.7)
            )
            wm_frame.fill.background()
            wm_frame.line.color.rgb = COLOR_GOLD
            wm_frame.line.width = Pt(1)

    def draw_glass_card(slide, left, top, width, height):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = COLOR_CARD
        card.fill.transparency = 0.12 # 12% transparent white container
        card.line.color.rgb = COLOR_GOLD
        card.line.width = Pt(1)
        return card

    def add_bullet_list(slide, left, top, width, height, bullets):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            
            parts = item.split(": ", 1)
            if len(parts) == 2:
                run1 = p.add_run()
                run1.text = parts[0] + ": "
                run1.font.name = 'Arial'
                run1.font.bold = True
                run1.font.size = Pt(13)
                run1.font.color.rgb = COLOR_TEXT
                
                run2 = p.add_run()
                run2.text = parts[1]
                run2.font.name = 'Arial'
                run2.font.size = Pt(13)
                run2.font.color.rgb = COLOR_MUTED
            else:
                run = p.add_run()
                run.text = item
                run.font.name = 'Arial'
                run.font.size = Pt(13)
                run.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 1: COVER SLIDE (IMMERSIVE BRIGHT BG)
    # ==========================================
    slide_1 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_1, img_hero)
    
    # Rounded Cover Card in Center
    draw_glass_card(slide_1, Inches(3.0), Inches(0.5), Inches(7.333), Inches(6.5))

    # Embed Logo on Cover Card
    if os.path.exists(logo_path):
        slide_1.shapes.add_picture(logo_path, Inches(5.666), Inches(0.8), Inches(2.0), Inches(2.0))
        logo_frame = slide_1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.666), Inches(0.8), Inches(2.0), Inches(2.0)
        )
        logo_frame.fill.background()
        logo_frame.line.color.rgb = COLOR_GOLD
        logo_frame.line.width = Pt(2)

    # Title
    logo_box = slide_1.shapes.add_textbox(Inches(3.2), Inches(3.0), Inches(6.933), Inches(0.8))
    logo_tf = logo_box.text_frame
    p_logo = logo_tf.paragraphs[0]
    p_logo.alignment = PP_ALIGN.CENTER
    run_logo = p_logo.add_run()
    run_logo.text = "DREAMHOME"
    run_logo.font.name = 'Georgia'
    run_logo.font.size = Pt(50)
    run_logo.font.bold = True
    run_logo.font.color.rgb = COLOR_GOLD

    # Subtitle
    sub_box = slide_1.shapes.add_textbox(Inches(3.2), Inches(3.8), Inches(6.933), Inches(0.4))
    sub_tf = sub_box.text_frame
    p_sub = sub_tf.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "THE PINNACLE OF LUXURY KITCHENS"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(11)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_BLUE
    
    # Divider Line
    line_shape = slide_1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5.666), Inches(4.35), Inches(2.0), Inches(0.02)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLOR_GOLD
    line_shape.line.color.rgb = COLOR_GOLD

    # Meta Details
    meta_box = slide_1.shapes.add_textbox(Inches(3.2), Inches(4.6), Inches(6.933), Inches(1.8))
    meta_tf = meta_box.text_frame
    
    p_meta1 = meta_tf.paragraphs[0]
    p_meta1.alignment = PP_ALIGN.CENTER
    p_meta1.text = "INVESTOR & COLLABORATION PROFILE"
    p_meta1.font.name = 'Arial'
    p_meta1.font.size = Pt(13)
    p_meta1.font.bold = True
    p_meta1.font.color.rgb = COLOR_TEXT
    p_meta1.space_after = Pt(10)

    p_meta2 = meta_tf.add_paragraph()
    p_meta2.alignment = PP_ALIGN.CENTER
    p_meta2.text = "# 301/1, B/W 16th & 17th Cross, Sampige Road, Malleshwaram, Bangalore - 560003"
    p_meta2.font.name = 'Arial'
    p_meta2.font.size = Pt(9)
    p_meta2.font.color.rgb = COLOR_MUTED
    p_meta2.space_after = Pt(4)

    p_meta3 = meta_tf.add_paragraph()
    p_meta3.alignment = PP_ALIGN.CENTER
    p_meta3.text = "Phone: 080-23344055, 9379599399 | Email: dreamhome.bengalure@gmail.com"
    p_meta3.font.name = 'Arial'
    p_meta3.font.size = Pt(9)
    p_meta3.font.color.rgb = COLOR_GOLD

    # ==========================================
    # SLIDE 2: THE SHOWROOM CONCEPT (MINIMAL WORDS)
    # ==========================================
    slide_2 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_2, img_display)
    add_slide_header(slide_2, "Experiential Showroom Hub")
    
    # Rounded card container
    draw_glass_card(slide_2, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.2))

    bullets_s2 = [
        "Strategic Hub: Malleshwaram prime retail node.",
        "Elite Demographic: Serving Bangalore high-net-worth.",
        "Live Test Labs: Active demonstrations in-store.",
        "High Yields: ₹8L average retail ticket."
    ]
    add_bullet_list(slide_2, Inches(1.1), Inches(2.0), Inches(5.4), Inches(3.8), bullets_s2)

    # Showroom Image on Right in Rounded Frame
    if os.path.exists(img_display):
        slide_2.shapes.add_picture(img_display, Inches(7.3), Inches(1.8), Inches(5.2), Inches(3.9))
        frame = slide_2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.45), Inches(1.95), Inches(5.2), Inches(3.9)
        )
        frame.fill.background()
        frame.line.color.rgb = COLOR_GOLD
        frame.line.width = Pt(1)
        slide_2.shapes._spTree.remove(frame._element)
        slide_2.shapes._spTree.insert(2, frame._element)

    # ==========================================
    # SLIDE 3: BRAND BOARD (BRIGHT ROUNDED CARDS)
    # ==========================================
    slide_3 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_3, img_hero)
    add_slide_header(slide_3, "Elite Brand Portfolio")

    # Main Card background (Rounded)
    draw_glass_card(slide_3, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.6))

    # Styled brands in corporate colors
    brands_styled = [
        ("FABER", "Italian Ventilation", RGBColor(227, 6, 19)),
        ("ELECTROLUX", "Swedish Eco Cooking", RGBColor(1, 30, 65)),
        ("SIEMENS", "German Smart Tech", RGBColor(0, 130, 138)),
        ("BOSCH", "German Laundry/Dish", RGBColor(227, 27, 35)),
        ("KAFF", "Premium Hobs/Hoods", RGBColor(30, 30, 30)),
        ("HÄFELE", "German Hardware", RGBColor(211, 1, 22)),
        ("HINDWARE", "Italian Collection", RGBColor(227, 30, 36)),
        ("CARYSIL", "Quartz Sinks", RGBColor(45, 45, 45)),
        ("A. O. SMITH", "Water Heaters/RO", RGBColor(0, 90, 156)),
        ("ZERO B", "Alkaline Purifiers", RGBColor(10, 94, 167)),
        ("RACOLD", "Thermal Heating", RGBColor(211, 47, 47)),
        ("DAIKIN", "Japanese HVAC", RGBColor(0, 160, 233))
    ]

    card_w = Inches(2.7)
    card_h = Inches(1.2)
    gap_x = Inches(0.2)
    gap_y = Inches(0.18)
    start_x = Inches(1.0)
    start_y = Inches(2.0)

    for idx, (bname, spec, color) in enumerate(brands_styled):
        row = idx // 4
        col = idx % 4
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        # Rounded brand cards
        card = slide_3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = COLOR_GOLD
        card.line.width = Pt(1)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.15)
        
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = bname
        p1.font.name = 'Georgia'
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = RGBColor(255, 255, 255)
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = spec
        p2.font.name = 'Arial'
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = RGBColor(230, 230, 230)

    # ==========================================
    # SLIDE 4: PRODUCT CATEGORIES (ROUNDED PILLARS)
    # ==========================================
    slide_4 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_4, img_display)
    add_slide_header(slide_4, "Curated Product Verticals")

    categories = [
        ("CULINARY", ["Designer Chimneys", "Built-in Hobs", "Built-in Ovens", "Microwave Ovens"]),
        ("REFRIGERATION", ["Built-in Cooling", "Luxury Refrigerators", "Beverage Cabinets"]),
        ("LAUNDRY & UTILITY", ["Built-in Dishwashers", "Washing Machines", "Heat Dryers"]),
        ("WATER & COMFORT", ["Quartz Composite Sinks", "Designer Faucets", "Water Purifiers", "Water Heaters"])
    ]

    col_w = Inches(2.7)
    start_x = Inches(0.8)
    start_y = Inches(1.8)
    gap = Inches(0.3)

    for idx, (title, items) in enumerate(categories):
        x = start_x + idx * (col_w + gap)
        
        box = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, col_w, Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.fill.transparency = 0.10
        box.line.color.rgb = COLOR_GOLD
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        
        p_title = tf.paragraphs[0]
        p_title.alignment = PP_ALIGN.CENTER
        p_title.text = title
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(11)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_GOLD
        p_title.space_after = Pt(15)

        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = "• " + item
            p_item.font.name = 'Arial'
            p_item.font.size = Pt(10.5)
            p_item.font.color.rgb = COLOR_TEXT
            p_item.space_after = Pt(8)

    # ==========================================
    # SLIDE 5: MARKET POTENTIAL (CONCISE METRICS)
    # ==========================================
    slide_5 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_5, img_hero)
    add_slide_header(slide_5, "Market Potential & Traction")

    # Left card
    draw_glass_card(slide_5, Inches(0.8), Inches(1.8), Inches(7.0), Inches(4.2))

    bullets_s5 = [
        "Built-in defaults: High modular kitchen upgrades.",
        "Market speed: Strong upscale home transaction growth.",
        "Allocated budgets: Growing premium appliance focus."
    ]
    add_bullet_list(slide_5, Inches(1.1), Inches(2.0), Inches(6.4), Inches(3.8), bullets_s5)

    # Stats Blocks on Right (Rounded)
    stats = [
        ("18%", "YoY CATEGORY GROWTH"),
        ("₹8L+", "AVERAGE SHOWROOM TICKET")
    ]
    for idx, (val, title) in enumerate(stats):
        y = Inches(1.8) + idx * Inches(2.2)
        box = slide_5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3), y, Inches(4.2), Inches(1.8))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CARD
        box.fill.transparency = 0.1
        box.line.color.rgb = COLOR_GOLD
        box.line.width = Pt(1)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        p1.text = val
        p1.font.name = 'Georgia'
        p1.font.size = Pt(36)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_GOLD
        
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.text = title
        p2.font.name = 'Arial'
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT

    # ==========================================
    # SLIDE 6: COLLABORATION MODELS (CONCISE TEXT)
    # ==========================================
    slide_6 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_6, img_display)
    add_slide_header(slide_6, "Collaboration Framework")

    # Center card (Rounded)
    draw_glass_card(slide_6, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))

    bullets_s6 = [
        "Retail showcase: Global brand dealership margins.",
        "Capital JVs: Scaling showroom cabinetry expansions.",
        "Designer network: Referral program for architects.",
        "Bulk builders: Developer pre-install kitchen packs."
    ]
    add_bullet_list(slide_6, Inches(1.2), Inches(2.1), Inches(11.0), Inches(3.9), bullets_s6)

    # ==========================================
    # SLIDE 7: SHOWROOM GALLERY (2 BRIGHT IMAGES ROUNDED)
    # ==========================================
    slide_7 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_7, None) # solid light background
    add_slide_header(slide_7, "Showroom Gallery")

    has_both = os.path.exists(img_hero) and os.path.exists(img_display)

    if has_both:
        # Two large images side-by-side
        slide_7.shapes.add_picture(img_hero, Inches(0.8), Inches(1.8), Inches(5.6), Inches(3.8))
        slide_7.shapes.add_picture(img_display, Inches(6.9), Inches(1.8), Inches(5.6), Inches(3.8))
        
        # Rounded Frames
        for x_coord in [Inches(0.8), Inches(6.9)]:
            frame = slide_7.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x_coord + Inches(0.08), Inches(1.88), Inches(5.6), Inches(3.8)
            )
            frame.fill.background()
            frame.line.color.rgb = COLOR_GOLD
            frame.line.width = Pt(1)
            slide_7.shapes._spTree.remove(frame._element)
            slide_7.shapes._spTree.insert(2, frame._element)

        tb = slide_7.shapes.add_textbox(Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.6))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = "Left: Sunlit Custom Showroom Concept  |  Right: Premium Built-In Hob & Chimney Display Workspace"
        p.font.name = 'Arial'
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 8: CONTACT DETAILS
    # ==========================================
    slide_8 = prs.slides.add_slide(blank_layout)
    apply_slide_background(slide_8, img_hero)
    add_slide_header(slide_8, "Initiate Collaboration")

    # Left card (Rounded)
    draw_glass_card(slide_8, Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.2))

    bullets_s8 = [
        "Address: # 301/1, B/W 16th & 17th Cross, Sampige Road, Malleshwaram, Bangalore - 560003.",
        "Store Contact: 080-23344055 | +91 93795 99399.",
        "Corporate Email: dreamhome.bengalure@gmail.com."
    ]
    add_bullet_list(slide_8, Inches(1.1), Inches(2.0), Inches(5.9), Inches(3.8), bullets_s8)

    # Right Card: CTA Callout (Rounded)
    card_cta = slide_8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.2))
    card_cta.fill.solid()
    card_cta.fill.fore_color.rgb = COLOR_CARD
    card_cta.fill.transparency = 0.12
    card_cta.line.color.rgb = COLOR_GOLD
    card_cta.line.width = Pt(1)

    tf_cta = card_cta.text_frame
    tf_cta.word_wrap = True
    tf_cta.margin_left = Inches(0.3)
    tf_cta.margin_right = Inches(0.3)
    tf_cta.margin_top = Inches(0.4)

    p_cta1 = tf_cta.paragraphs[0]
    p_cta1.alignment = PP_ALIGN.CENTER
    p_cta1.text = "LET'S BUILD TOGETHER"
    p_cta1.font.name = 'Arial'
    p_cta1.font.size = Pt(14)
    p_cta1.font.bold = True
    p_cta1.font.color.rgb = COLOR_GOLD
    p_cta1.space_after = Pt(20)

    p_cta2 = tf_cta.add_paragraph()
    p_cta2.alignment = PP_ALIGN.LEFT
    p_cta2.text = "We welcome brands, designers, and developers to schedule a private showroom walkthrough or margin review."
    p_cta2.font.name = 'Arial'
    p_cta2.font.size = Pt(11)
    p_cta2.font.color.rgb = COLOR_TEXT
    p_cta2.space_after = Pt(14)

    p_cta3 = tf_cta.add_paragraph()
    p_cta3.alignment = PP_ALIGN.LEFT
    p_cta3.text = "Connect with us to explore floorplans, margin structures, and dealer frameworks."
    p_cta3.font.name = 'Arial'
    p_cta3.font.size = Pt(11)
    p_cta3.font.color.rgb = COLOR_MUTED

    # 3. Save
    filename = "dreamhome_profile.pptx"
    fallback_filename = "dreamhome_profile_v2.pptx"
    try:
        prs.save(filename)
        print(f"PowerPoint presentation '{filename}' generated successfully!")
    except PermissionError:
        print(f"Warning: '{filename}' is locked (likely open in PowerPoint). Saving to '{fallback_filename}' instead.")
        prs.save(fallback_filename)
        print(f"PowerPoint presentation '{fallback_filename}' generated successfully!")

if __name__ == "__main__":
    create_presentation()
